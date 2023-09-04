import collections 
import collections.abc
from pptx import Presentation
import wikipedia
import requests # request img from web
import uuid
import os
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.text.text import MSO_AUTO_SIZE
from pptx.text.text import TextFrame
from pptx.text.text import Font
from pptx.util import Inches
#Working
from pptx.util import Pt 
import pptx
import re
from random import randint
from point_plane_ui import Ui_mainWindow
import sys
import webbrowser
from PIL import Image
from PyQt5.QtCore import QObject, QThread, pyqtSignal


import os
import subprocess


import threading

import time

print(pptx.__file__)

#!!!!!!!!!!!!!!!!!!remove the "_rc" from the end of point_plane_res and in the ui file too
import point_plane_res

from PyQt5.QtWidgets import (

    QApplication, QDialog, QMainWindow, QMessageBox

)
from PyQt5.uic import loadUi




    














def transliterate(string):
    from unidecode import unidecode

    if not isinstance(string, bytes):
        string = u''.join(string)

    return unidecode(string)


prs = Presentation('/home/abris/.local/lib/python3.10/site-packages/pptx/templates/default.pptx')

#bg1, bg2, titles, text
color_pallette= [RGBColor(96, 153, 102), RGBColor(157, 192, 139), RGBColor(64, 81, 59), RGBColor(44, 61, 39)]
color_pallette_ = [RGBColor(232, 170, 66), RGBColor(229, 124, 35), RGBColor(52, 134, 150), RGBColor(2, 84, 100)]

headers = {'User-Agent': 'PointPlaneImageBot/0.1 (https://www.instagram.com/sabris_abris/; suhai.abris@gmail.com)'}

def downloadimages(query, dir):




    if query.endswith(".svg") or query.endswith(".SVG"):

        print("query:" + query)
        
        file_name = re.sub(r'[^a-zA-Z0-9\s]', '', query.rsplit('/', 1)[-1]).replace("png", ".png").replace("jpg", ".jpg").replace("svg", ".svg").replace("JPG", ".JPG").replace("jpeg", ".jpeg").replace("PNG", ".PNG").replace("JPEG", ".JPEG").replace("SVG", ".SVG")

        print(query)

        res = requests.get(query, headers=headers)

        if res.status_code == 200:
            with open(os.path.join(dir, file_name),'wb') as f:
                f.write(res.content)
            print('Image sucessfully Downloaded: ',file_name)
            return "403"
        else:
            print('Image Couldn\'t be retrieved' + str(res.status_code))
            return "403"

    
    if query.endswith(".png") or query.endswith(".jpg") or query.endswith(".jpeg") or query.endswith(".JPG") or query.endswith(".PNG"):
        print("query:" + query)
        
        file_name = re.sub(r'[^a-zA-Z0-9\s]', '', query.rsplit('/', 1)[-1]).replace("png", ".png").replace("jpg", ".jpg").replace("svg", ".svg").replace("JPG", ".JPG").replace("jpeg", ".jpeg").replace("PNG", ".PNG").replace("JPEG", ".JPEG").replace("SVG", ".SVG")

        print(query)

        res = requests.get(query, headers=headers)

        if res.status_code == 200:
            with open(os.path.join(dir, file_name),'wb') as f:
                f.write(res.content)
            print('Image sucessfully Downloaded: ',file_name)
            return os.path.join(dir, file_name)
        else:
            print('Image Couldn\'t be retrieved' + str(res.status_code))
            return "403"
    else:
        return "403"
    
    

    





def split_text_into_sentences(text):
    # Split text into sentences using regular expressions
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)
    return sentences


def combine_sentences(sentences, max_characters):
    combined_sentences = []
    character_count = 0

    for sentence in sentences:
        if character_count + len(sentence) <= max_characters:
            combined_sentences.append(sentence)
            character_count += len(sentence)
        else:
            break

    # If only 4 sentences are added and total character count still exceeds the limit, remove the last sentence
    if len(combined_sentences) == 5 and character_count > max_characters:
        combined_sentences.pop()

    return combined_sentences



def alert(ui, txt, title):

    ui.fog0.show()
    ui.fog2.show()
    ui.fog3.show()
    ui.fog4.show()
    ui.fog5.show()
    ui.okBtn.show()

    ui.fog3.setText(title)
    ui.fog5.setText(txt)


class Worker(QObject):

    finished = pyqtSignal()


    

    def gen_slides(self, pag, lang, ui):


        
        
        wikipedia.set_lang(lang)


        try:
            page = wikipedia.page(pag, auto_suggest=False)

        except wikipedia.DisambiguationError as e:
            print(e.options)

            alert(ui, str(e.options), "NOT FOUND, PLEASE CHOOSE AN OPTION")

            return
        

        
            

        
        ui.progressBar.setValue(0)

        if not os.path.exists(page.title):
            os.mkdir(page.title)





        print(page.images)

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title




        filll = slide.background.fill


        filll.gradient()
                
        filll.gradient_angle = 90


        # 1 - 6 ig vannak a themek
        filll.gradient_stops[0].color.rgb = color_pallette[0]
        filll.gradient_stops[1].color.rgb = color_pallette[1]

        filll.gradient_stops[0].position = 0.3



        title.text = page.title
        title.text_frame.paragraphs[0].font.name = "Ubuntu"
        title.text_frame.paragraphs[0].font.color.rgb = color_pallette[2]

        print(page.sections)

        

        sec = page.sections

        try:
            sec.remove("Kapcsolódó cikkek")
        except:
            pass
        try:
            sec.remove("Külső hivatkozások")
        except:
            pass
        try:
            sec.remove("Kapcsolódó szócikkek")
        except:
            pass
        try:
            sec.remove("Jegyzetek")
        except:
            pass
        try:
            sec.remove("Lásd még")
        except:
            pass
        try:
            sec.remove("Fordítás")
        except:
            pass
        try:
            sec.remove("További információk")
        except:
            pass



        dwlimages = []


        for url in page.images:

            ui.progressBar.setValue(int((100 / len(page.images)) * page.images.index(url)))

            pc = downloadimages(url, page.title)

            if pc != "403":
                dwlimages.append(pc)





        for i in sec:

            sectio = page.section(i)


            
            if sectio != "":

                result = list(filter(lambda x : x != '', sectio.split('\n\n')))

                cnt_slide = prs.slides.add_slide(prs.slide_layouts[7])

                fill = cnt_slide.background.fill

                fill.gradient()
                
                fill.gradient_angle = 45


                # 1 - 6 ig vannak a themek
                fill.gradient_stops[0].color.rgb = color_pallette[0]
                fill.gradient_stops[1].color.rgb = color_pallette[1]
                fill.gradient_stops[0].position = 0.3


                

                cnt_slide.shapes.title.text = i




                max_characters = 450

                sentences = split_text_into_sentences(result[0])
                combined_sentences = combine_sentences(sentences, max_characters)

                cnt_slide.placeholders[2].text = ' '.join(combined_sentences)


                font = cnt_slide.placeholders[2].text_frame.paragraphs[0].font

                tfont = cnt_slide.placeholders[0].text_frame.paragraphs[0].font

                tfont.size = Pt(30)
                tfont.name = "Ubuntu"
                tfont.color.rgb = color_pallette[2]

                font.size = Pt(15)
                font.name = "Ubuntu"
                font.color.rgb = color_pallette[3]

                
                pim = page.images


                for im in pim:
                    if i.lower().replace(" ", "") in im.lower().replace('_', ''):
                        print(i.lower().replace('_', '').replace(" ", ""))
                        print(im.lower().replace('_', '').replace(" ", ""))

                        image = downloadimages(im, page.title)
                    else:
                        image = "403"

                left = top = 0

                text_box_width = Inches(6.1)  # Width of the text box on the left
                total_slide_width = Inches(10)  # Total width of the slide

                # Calculate the left position for the image
                left_position = total_slide_width - text_box_width  # Position the image just after the text box
                top_position = Inches(1.75)
                image_width = Inches(6)  # Adjust this as needed
                image_height = Inches(4)

                if image != "403":
                    pic = cnt_slide.shapes.add_picture(image, left_position, top_position, image_width, image_height)

                
                else:
                    if len(dwlimages) > 0:

                        if Image.open(dwlimages[0]).format != "MPO":
                            
                            print(Image.open(dwlimages[0]).format)

                            pic = cnt_slide.shapes.add_picture(dwlimages[0], left_position, top_position, image_width, image_height)
                        else:
                            print("MMMMMMPPPPPOOOOOOO")


                        dwlimages.pop(0)

                        print(dwlimages)


        prs.save(f'{str(page.title)}.pptx')

        

        print("saved")


        ui.progressBar.setValue(0)

        self.finished.emit()





class Window(QMainWindow, Ui_mainWindow):


    def open_dir(self, dir_path, fold):

        try:
            subprocess.run(["open", dir_path], check=True)  # On macOS
        except FileNotFoundError:
            try:
                subprocess.run(["xdg-open", dir_path], check=True)  # On Linux
            except FileNotFoundError:
                try:
                    subprocess.run(["start", dir_path], check=True, shell=True)  # On Windows
                except FileNotFoundError:
                    print("Unable to open the file. Please check the file path or your operating system.")


        if sys.platform=='win32':
            subprocess.Popen(['start', os.path.join(os.getcwd(), fold)], shell= True)

        elif sys.platform=='darwin':
            subprocess.Popen(['open', os.path.join(os.getcwd(), fold)])

        else:
            try:
                subprocess.Popen(['xdg-open', os.path.join(os.getcwd(), fold)])
            except OSError:
                print("Your crappy os is not fully supported")


    


    def startG(self):


        if len(self.titleInput.text()) > 0:

            self.startBtn.setEnabled(False)

            self.startBtn.setText("LOADING...")

            

            print("Started generating Slide Show from input: " + self.titleInput.text())

            # Step 2: Create a QThread object
            self.thread = QThread()
            # Step 3: Create a worker object

            
            title_input_text = self.titleInput.text()
            current_lang = self.setLang.currentText()


            self.worker = Worker()


            # Step 4: Move worker to the thread
            self.worker.moveToThread(self.thread)

            # Step 5: Connect signals and slots
            self.thread.started.connect(lambda: self.worker.gen_slides(self.titleInput.text(), self.setLang.currentText(), self))



            self.worker.finished.connect(self.thread.quit)
            self.worker.finished.connect(self.worker.deleteLater)
            self.thread.finished.connect(self.thread.deleteLater)

            # Step 6: Start the thread
            self.thread.start()

            # Final resets
            self.thread.finished.connect(lambda: self.reset(wikipedia.page(self.titleInput.text(), auto_suggest=False).title))


        else:
            print("No Input")


    def reset(self, about):

        self.startBtn.setEnabled(True)

        self.startBtn.setText("START")

        

        


        self.fog0.show()
        self.fog2.show()

        self.fog3.show()
        self.fog3.setText("INFO")

        self.fog4.show()

        self.fog5.show()
        self.fog5.setText(f"Successfully generated ppt. You can use the images in the folder called {about}. You can manually add them into the presentation.")

        self.okBtn.show()
        self.okBtn.setText("OK")

        self.open_dir(f'{str(about)}.pptx', about)


    def __init__(self, parent=None):

        super().__init__(parent)

        self.setupUi(self)

        self.connectSignalsSlots()

        self.progressBar.setValue(0)

    def setLanguage(self):
        print(self.setLang.currentText())
        wikipedia.set_lang(self.setLang.currentText())



    def supportSite(self):
        print("+1 süti")


    def connectSignalsSlots(self):



        self.startBtn.clicked.connect(self.startG)
        self.supportBtn.clicked.connect(self.supportSite)
        self.setLang.currentIndexChanged.connect(self.setLanguage)

        progressB = self.progressBar



if __name__ == "__main__":

    app = QApplication(sys.argv)

    win = Window()

    win.show()

    sys.exit(app.exec())